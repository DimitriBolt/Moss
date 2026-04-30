from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT
from reportlab.lib.utils import ImageReader
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.pdfgen import canvas
from reportlab.platypus import Paragraph
from reportlab.lib.styles import ParagraphStyle
from xml.sax.saxutils import escape


ROOT = Path(__file__).resolve().parents[1]
POSTER = ROOT / "Poster"
IMG = ROOT / "img"

TEMPLATE = POSTER / "42x60 Poster Template - FINAL_2023 (1).pptx"
OUT = POSTER / "moss_spontaneous_colonization_poster.pptx"
PDF_OUT = POSTER / "moss_spontaneous_colonization_poster.pdf"

UA_MATH_LOGO = POSTER / "ua_math_logo.png"
BIOSPHERE2_LOGO = POSTER / "biosphere2_logo.png"

CENTER_IMAGES = [
    (IMG / "LEO_2022_11_22_cropped_rectified_rectangle.png", "Rectified LEO1 mosaic", "5-camera mosaic registered to bath coordinates"),
    (IMG / "persistent_moss_map.png", "Persistent moss footprint", "Cells positive on >=2 of 3 analysis dates"),
    (IMG / "moss_cross_sections.png", "Topographic cross-sections", "Trough minima are repeatedly moss-poor"),
]

ARTICLE_RESPONSE_CURVES = POSTER / "moss_response_curves.png"
ARTICLE_CALIBRATION = POSTER / "stable_moss_calibration.png"


COLORS = {
    "sage": RGBColor(160, 168, 139),
    "sage_dark": RGBColor(119, 128, 100),
    "taupe": RGBColor(137, 127, 99),
    "taupe_dark": RGBColor(92, 82, 63),
    "cream": RGBColor(246, 247, 242),
    "card": RGBColor(255, 255, 255),
    "ink": RGBColor(24, 25, 23),
    "muted": RGBColor(73, 76, 70),
    "blue": RGBColor(0, 47, 108),
    "red": RGBColor(171, 5, 32),
    "line": RGBColor(210, 213, 202),
    "green": RGBColor(67, 119, 82),
}


def I(value: float):
    return Inches(value)


def clear_slide(slide):
    for shape in list(slide.shapes):
        shape._element.getparent().remove(shape._element)


def delete_extra_slides(prs, keep_count=1):
    sld_id_lst = prs.slides._sldIdLst
    for sld_id in list(sld_id_lst)[keep_count:]:
        prs.part.drop_rel(sld_id.rId)
        sld_id_lst.remove(sld_id)


def add_rect(slide, x, y, w, h, fill, line=None, radius=False, transparency=0):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, I(x), I(y), I(w), I(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.fill.transparency = transparency
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.6)
    return shp


def add_text(slide, x, y, w, h, text, size=11, color=None, bold=False,
             italic=False, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP,
             font="DejaVu Sans", margin=0.06, line_spacing=1.0):
    box = slide.shapes.add_textbox(I(x), I(y), I(w), I(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = I(margin)
    tf.margin_right = I(margin)
    tf.margin_top = I(margin)
    tf.margin_bottom = I(margin)
    tf.vertical_anchor = valign

    lines = text.split("\n")
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.line_spacing = line_spacing
        for run in p.runs:
            run.font.name = font
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = color or COLORS["ink"]
    return box


def add_mixed_line(slide, x, y, w, h, pieces, size=10.5, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(I(x), I(y), I(w), I(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = I(0.04)
    tf.margin_right = I(0.04)
    tf.margin_top = I(0.02)
    tf.margin_bottom = I(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    for text, bold, color in pieces:
        r = p.add_run()
        r.text = text
        r.font.name = "DejaVu Sans"
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return box


def section_header(slide, x, y, w, title, color=None):
    add_rect(slide, x, y, w, 0.38, color or COLORS["taupe"], line=None)
    add_text(
        slide, x + 0.08, y + 0.035, w - 0.16, 0.28, title,
        size=12.5, color=RGBColor(255, 255, 255), bold=True,
        align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0.0
    )


def card(slide, x, y, w, h, title, title_color=None):
    add_rect(slide, x, y, w, h, COLORS["card"], line=COLORS["line"])
    section_header(slide, x, y, w, title, title_color)
    return (x + 0.18, y + 0.52, w - 0.36, h - 0.68)


def add_bullets(slide, x, y, w, h, items, size=10.6, color=None, gap=0.19):
    body = "\n".join(f"- {item}" for item in items)
    return add_text(slide, x, y, w, h, body, size=size, color=color or COLORS["ink"], line_spacing=1.0, margin=0.03)


def add_picture_fit(slide, path, x, y, w, h, border=True, fill=RGBColor(255, 255, 255)):
    add_rect(slide, x, y, w, h, fill, line=COLORS["line"] if border else None)
    with Image.open(path) as img:
        iw, ih = img.size
    aspect = iw / ih
    box_aspect = w / h
    if aspect >= box_aspect:
        pw = w
        ph = w / aspect
    else:
        ph = h
        pw = h * aspect
    pic = slide.shapes.add_picture(str(path), I(x + (w - pw) / 2), I(y + (h - ph) / 2), width=I(pw), height=I(ph))
    return pic


def add_stat(slide, x, y, w, h, value, label, accent=None):
    add_rect(slide, x, y, w, h, COLORS["cream"], line=COLORS["line"])
    add_text(slide, x + 0.05, y + 0.07, w - 0.1, 0.32, value, size=15, color=accent or COLORS["green"], bold=True,
             align=PP_ALIGN.CENTER, margin=0.0)
    add_text(slide, x + 0.08, y + 0.45, w - 0.16, h - 0.48, label, size=7.5, color=COLORS["muted"],
             align=PP_ALIGN.CENTER, margin=0.0)


def build():
    prs = Presentation(str(TEMPLATE))
    delete_extra_slides(prs, keep_count=1)
    slide = prs.slides[0]
    clear_slide(slide)
    prs.slide_width = I(30)
    prs.slide_height = I(21)

    # Background and header
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["cream"]
    add_rect(slide, 0, 0, 30, 2.42, COLORS["sage"], line=None)
    add_rect(slide, 0, 2.42, 30, 0.08, COLORS["taupe_dark"], line=None)

    add_picture_fit(slide, UA_MATH_LOGO, 0.55, 0.48, 4.85, 1.18, border=False, fill=COLORS["sage"])
    add_picture_fit(slide, BIOSPHERE2_LOGO, 25.05, 0.68, 4.35, 0.85, border=False, fill=COLORS["sage"])

    title = (
        "Spontaneous Moss Colonization on Crushed Basalt:\n"
        "Imaging Pipeline, Color Reliability, and Topographic Controls on Persistence"
    )
    add_text(slide, 5.25, 0.28, 19.5, 1.15, title, size=23.5, color=RGBColor(0, 0, 0),
             bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0.0, line_spacing=0.9)
    add_text(slide, 7.4, 1.45, 15.2, 0.35, "Freddy Aguilar, Sophia Bui, Ava Mohr", size=17,
             color=RGBColor(0, 0, 0), align=PP_ALIGN.CENTER, margin=0.0)
    add_text(slide, 7.4, 1.84, 15.2, 0.32, "The University of Arizona, Department of Mathematics",
             size=11.8, color=RGBColor(0, 0, 0), italic=True, align=PP_ALIGN.CENTER, margin=0.0)

    # Column geometry
    left_x, left_w = 0.5, 6.35
    center_x, center_w = 7.15, 12.05
    right_x, right_w = 19.55, 9.95
    top_y = 2.8

    # Left column
    x, y, w, h = card(slide, left_x, top_y, left_w, 3.05, "Why This System?")
    add_bullets(slide, x, y, w, h, [
        "Biosphere 2's LEO Observatory is a controlled 30 m x 30 m hillslope built from crushed basalt.",
        "The substrate is a Mars-analog, soil-poor surface, yet moss and biological soil crusts colonized it without planting.",
        "The scientific question is whether this colonization is random or structured by hillslope geometry."
    ], size=9.8)

    x, y, w, h = card(slide, left_x, 6.1, left_w, 2.15, "Research Questions", COLORS["sage_dark"])
    add_text(slide, x + 0.05, y, w - 0.1, h, "Q1. Can overhead images reliably map moss area?\n\nQ2. Do topographic features predict where moss persists?",
             size=11.4, color=COLORS["ink"], bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    x, y, w, h = card(slide, left_x, 8.5, left_w, 2.75, "Experimental Setup")
    add_bullets(slide, x, y, w, h, [
        "Sensor-rich laboratory system: 1,800+ devices, 16 sensor columns, and 5 cameras per slope.",
        "LiDAR DEM at 0.1 m resolution supports geometry-based analysis.",
        "Image archive includes RGB and RGB plus near-infrared imagery, filtered by quality."
    ], size=9.8)

    x, y, w, h = card(slide, left_x, 11.5, left_w, 2.65, "Previous Work", COLORS["sage_dark"])
    add_bullets(slide, x, y, w, h, [
        "Klausmeier (1999): slope water redistribution can generate vegetation stripes and spots.",
        "Rietkerk et al. (2002): plant-water-infiltration feedbacks produce non-uniform biological structure.",
        "This study tests that patterning logic on a controlled basalt hillslope."
    ], size=9.4)

    x, y, w, h = card(slide, left_x, 14.4, left_w, 3.15, "Data and Time Window")
    add_bullets(slide, x, y, w, h, [
        "Analysis dates: Nov 3, Nov 22, and Dec 7, 2022.",
        "Primary response: persistent moss mask from repeated positive detections.",
        "Topographic covariates: lateral distance, downstream position, trough depth, flow proxy, wall proximity, and soil thickness.",
        "Previous poster time frame: 2022-10-08 to 2024-01-01."
    ], size=9.35)

    x, y, w, h = card(slide, left_x, 17.8, left_w, 2.25, "Poster Story", COLORS["taupe_dark"])
    add_text(slide, x, y, w, h,
             "Act 1: a spontaneous Mars-analog colonization problem.\n"
             "Act 2: reliable imaging requires multi-date persistence, not single-date color.\n"
             "Act 3: topography predicts persistence through a bounded moisture optimum.",
             size=9.35, color=COLORS["ink"], margin=0.03)

    # Center column
    x, y, w, h = card(slide, center_x, top_y, center_w, 2.8, "Methodology: Image-to-Moss Pipeline")
    stage_w = (w - 0.42) / 3
    stages = [
        ("01", "Calibration", "Radial lens correction and per-camera rectification."),
        ("02", "Stitching", "Five corrected frames merged into one bath-coordinate mosaic."),
        ("03", "Persistence Proxy", "Broad RGB/HSV moss proxy retained only when positive on >=2 of 3 dates."),
    ]
    for i, (num, head, desc) in enumerate(stages):
        sx = x + i * (stage_w + 0.21)
        add_rect(slide, sx, y + 0.02, stage_w, h - 0.06, COLORS["cream"], line=COLORS["line"])
        add_text(slide, sx + 0.08, y + 0.12, 0.48, 0.34, num, size=13, color=COLORS["red"], bold=True, margin=0.0)
        add_text(slide, sx + 0.62, y + 0.12, stage_w - 0.72, 0.34, head, size=11.2, color=COLORS["blue"], bold=True, margin=0.0)
        add_text(slide, sx + 0.1, y + 0.56, stage_w - 0.2, h - 0.66, desc, size=8.9, color=COLORS["ink"], margin=0.0)

    x, y, w, h = card(slide, center_x, 5.9, center_w, 9.55, "Central Evidence Panels", COLORS["sage_dark"])
    img_y = y + 0.02
    center_slots = [
        (CENTER_IMAGES[0], x, img_y, 3.35, 7.65, img_y + 7.82),
        (CENTER_IMAGES[1], x + 3.58, img_y, 3.35, 5.15, img_y + 5.32),
        (CENTER_IMAGES[2], x + 7.16, img_y, 4.53, 5.55, img_y + 5.72),
    ]
    for (path, head, desc), sx, sy, sw, sh, cy in center_slots:
        add_picture_fit(slide, path, sx, sy, sw, sh)
        add_text(slide, sx, cy, sw, 0.28, head, size=9.4, color=COLORS["blue"], bold=True,
                 align=PP_ALIGN.CENTER, margin=0.0)
        add_text(slide, sx + 0.05, cy + 0.31, sw - 0.1, 0.52, desc, size=7.4, color=COLORS["muted"],
                 align=PP_ALIGN.CENTER, margin=0.0)

    add_rect(slide, x, y + 8.72, w, 0.62, COLORS["taupe_dark"], line=None)
    add_text(slide, x + 0.16, y + 8.8, w - 0.32, 0.45,
             "Stable moss is concentrated on off-center shoulder zones, not in the deepest drainage corridor.",
             size=10.9, color=RGBColor(255, 255, 255), bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0.0)

    x, y, w, h = card(slide, center_x, 15.75, center_w, 4.3, "Findings 01-02: Persistence Beats Color")
    add_text(slide, x, y, w * 0.58, h - 0.1,
             "Finding 01 - Moss is spatially persistent.\n"
             "The same regions remain moss-positive in November and December. Later greening is mostly color intensification in already occupied regions.\n\n"
             "Finding 02 - Color thresholding is unreliable.\n"
             "Hydration and irrigation can shift moss from yellow-green to vivid green within hours, so single-date green area can mimic growth.",
             size=9.25, color=COLORS["ink"], margin=0.03)
    stat_x = x + w * 0.61
    add_stat(slide, stat_x, y + 0.04, 1.55, 1.05, "26.6%", "stable on >=2 dates")
    add_stat(slide, stat_x + 1.75, y + 0.04, 1.55, 1.05, "7.1%", "strong on all 3 dates")
    add_stat(slide, stat_x, y + 1.33, 1.55, 1.05, "3.7 m", "median lateral offset")
    add_stat(slide, stat_x + 1.75, y + 1.33, 1.55, 1.05, "0.088 m", "median trough depth")
    add_text(slide, stat_x, y + 2.66, 3.28, 0.82,
             "Interpretation: the response variable is persistence, not raw greenness.",
             size=8.5, color=COLORS["taupe_dark"], bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0.02)

    # Right column
    x, y, w, h = card(slide, right_x, top_y, right_w, 2.9, "Finding 03: Topographic Control", COLORS["taupe_dark"])
    add_rect(slide, x, y + 0.02, w, 0.72, COLORS["blue"], line=None)
    add_text(slide, x + 0.12, y + 0.12, w - 0.24, 0.5,
             "The wettest corridor is not the most favorable habitat.",
             size=12.5, color=RGBColor(255, 255, 255), bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0.0)
    add_bullets(slide, x, y + 0.92, w, h - 1.0, [
        "Persistent moss increases away from the centerline.",
        "Trough-depth response is non-monotone, peaking near an intermediate depth.",
        "Deep, high-flow drainage corridors behave as exclusion zones."
    ], size=9.4)

    x, y, w, h = card(slide, right_x, 6.0, right_w, 5.05, "Quantitative Evidence")
    add_picture_fit(slide, ARTICLE_RESPONSE_CURVES, x, y + 0.02, w * 0.54, h - 0.34)
    add_picture_fit(slide, ARTICLE_CALIBRATION, x + w * 0.57, y + 0.02, w * 0.43, h - 0.34)
    add_text(slide, x, y + h - 0.27, w, 0.22,
             "Response curves and suitability-decile calibration from the technical note.",
             size=7.4, color=COLORS["muted"], align=PP_ALIGN.CENTER, margin=0.0)

    x, y, w, h = card(slide, right_x, 11.35, right_w, 3.55, "Logistic Suitability Ranking Model", COLORS["sage_dark"])
    eq = "logit P(M=1) = beta0 + beta_w W - beta_f F - gamma(d-d0)^2 + beta_t T"
    add_rect(slide, x, y + 0.04, w, 0.57, COLORS["cream"], line=COLORS["line"])
    add_text(slide, x + 0.08, y + 0.16, w - 0.16, 0.32, eq, size=8.9, color=COLORS["ink"],
             bold=True, align=PP_ALIGN.CENTER, margin=0.0)
    add_text(slide, x, y + 0.82, w * 0.51, 1.74,
             "Predictors\n"
             "W: wall proximity (+)\n"
             "F: log(1+flow) (-)\n"
             "d: trough depth, optimum d0=0.407 m\n"
             "T: soil thickness (+)",
             size=8.7, color=COLORS["ink"], margin=0.02)
    add_text(slide, x + w * 0.55, y + 0.82, w * 0.42, 1.74,
             "Performance\n"
             "Training AUC: 0.724\n"
             "Spatial CV AUC: 0.657\n\n"
             "Suitability deciles rank stable-moss fraction from about 7% to 48%.",
             size=8.7, color=COLORS["ink"], margin=0.02)
    add_text(slide, x, y + 2.72, w, 0.34,
             "Use as first-order ranking, not a complete mechanistic law.",
             size=8.6, color=COLORS["taupe_dark"], bold=True, align=PP_ALIGN.CENTER, margin=0.0)

    x, y, w, h = card(slide, right_x, 15.2, right_w, 2.35, "Why the Trough May Suppress Moss")
    mechanisms = [
        "excess saturation and low oxygen",
        "surface flushing of spores and fines",
        "sediment redistribution",
        "mechanical surface instability",
        "moisture-window limitation"
    ]
    add_bullets(slide, x, y, w, h, mechanisms, size=8.8)

    x, y, w, h = card(slide, right_x, 17.85, right_w, 2.2, "Conclusions and Next Steps", COLORS["taupe_dark"])
    add_text(slide, x, y, w * 0.54, h - 0.05,
             "Q1: Yes - use multi-date persistence to suppress hydration artifacts.\n\n"
             "Q2: Yes - topography ranks persistence, with shoulder zones favored over active drainage corridors.\n\n"
             "Implication: pioneer colonizers occupy off-channel shoulders on Mars-analog basalt.",
             size=8.85, color=COLORS["ink"], margin=0.02)
    add_text(slide, x + w * 0.58, y, w * 0.4, h - 0.05,
             "Next measurements\n"
             "- soil moisture time series\n"
             "- saturation duration\n"
             "- tracer residence time\n"
             "- CO2 flux vs moss density",
             size=8.65, color=COLORS["ink"], margin=0.02)

    # Footer
    add_text(slide, 0.55, 20.35, 13.4, 0.28,
             "Scientific source: moss_technical_note_scientific_revision.pdf; poster narrative: Findings-content.md.md",
             size=6.8, color=COLORS["muted"], margin=0.0)
    add_text(slide, 20.0, 20.35, 9.45, 0.28,
             "References: Klausmeier 1999; Rietkerk et al. 2002; LEO Observatory / Biosphere 2 context.",
             size=6.8, color=COLORS["muted"], align=PP_ALIGN.RIGHT, margin=0.0)

    prs.save(str(OUT))
    print(OUT)
    build_pdf()


def rl_color(rgb):
    return colors.Color(rgb[0] / 255, rgb[1] / 255, rgb[2] / 255)


def register_pdf_fonts():
    font_dir = Path("/usr/share/fonts/truetype/dejavu")
    pdfmetrics.registerFont(TTFont("DejaVuSans", str(font_dir / "DejaVuSans.ttf")))
    pdfmetrics.registerFont(TTFont("DejaVuSans-Bold", str(font_dir / "DejaVuSans-Bold.ttf")))
    pdfmetrics.registerFont(TTFont("DejaVuSans-Oblique", str(font_dir / "DejaVuSans-Oblique.ttf")))


class PDFDeck:
    def __init__(self, path):
        register_pdf_fonts()
        self.page_w = 30 * 72
        self.page_h = 21 * 72
        self.c = canvas.Canvas(str(path), pagesize=(self.page_w, self.page_h))

    def px(self, x):
        return x * 72

    def py(self, y, h=0):
        return self.page_h - (y + h) * 72

    def rect(self, x, y, w, h, fill, line=None, stroke_width=0.6):
        self.c.setFillColor(rl_color(fill))
        if line is None:
            self.c.setStrokeColor(rl_color(fill))
            self.c.setLineWidth(0)
        else:
            self.c.setStrokeColor(rl_color(line))
            self.c.setLineWidth(stroke_width)
        self.c.rect(self.px(x), self.py(y, h), self.px(w), self.px(h), fill=1, stroke=0 if line is None else 1)

    def text(self, x, y, w, h, text, size=10, color=None, bold=False, italic=False,
             align="left", valign="top", leading=None):
        font = "DejaVuSans-Bold" if bold else "DejaVuSans-Oblique" if italic else "DejaVuSans"
        alignment = {"left": TA_LEFT, "center": TA_CENTER, "right": TA_RIGHT}[align]
        style = ParagraphStyle(
            "p",
            fontName=font,
            fontSize=size,
            leading=leading or size * 1.18,
            textColor=rl_color(color or COLORS["ink"]),
            alignment=alignment,
            spaceAfter=0,
            spaceBefore=0,
        )
        safe = escape(text).replace("\n", "<br/>")
        para = Paragraph(safe, style)
        max_w, max_h = self.px(w), self.px(h)
        used_w, used_h = para.wrap(max_w, max_h)
        draw_x = self.px(x)
        if valign == "middle":
            draw_y = self.py(y, h) + max(0, (max_h - used_h) / 2)
        elif valign == "bottom":
            draw_y = self.py(y, h)
        else:
            draw_y = self.page_h - self.px(y) - used_h
        para.drawOn(self.c, draw_x, draw_y)

    def header(self, x, y, w, title, fill=None):
        self.rect(x, y, w, 0.38, fill or COLORS["taupe"])
        self.text(x + 0.08, y + 0.055, w - 0.16, 0.26, title, size=12.4,
                  color=RGBColor(255, 255, 255), bold=True, align="center", valign="middle")

    def card(self, x, y, w, h, title, fill=None):
        self.rect(x, y, w, h, COLORS["card"], COLORS["line"])
        self.header(x, y, w, title, fill)
        return x + 0.18, y + 0.52, w - 0.36, h - 0.68

    def bullets(self, x, y, w, h, items, size=9.5):
        self.text(x, y, w, h, "\n".join(f"- {item}" for item in items), size=size)

    def image_fit(self, path, x, y, w, h, border=True, fill=RGBColor(255, 255, 255)):
        self.rect(x, y, w, h, fill, COLORS["line"] if border else None)
        with Image.open(path) as img:
            iw, ih = img.size
        aspect = iw / ih
        box_aspect = w / h
        if aspect >= box_aspect:
            pw, ph = w, w / aspect
        else:
            ph, pw = h, h * aspect
        dx = x + (w - pw) / 2
        dy = y + (h - ph) / 2
        self.c.drawImage(ImageReader(str(path)), self.px(dx), self.py(dy, ph), self.px(pw), self.px(ph), mask="auto")

    def stat(self, x, y, w, h, value, label, accent=None):
        self.rect(x, y, w, h, COLORS["cream"], COLORS["line"])
        self.text(x + 0.05, y + 0.07, w - 0.1, 0.32, value, size=15,
                  color=accent or COLORS["green"], bold=True, align="center")
        self.text(x + 0.08, y + 0.45, w - 0.16, h - 0.48, label, size=7.5,
                  color=COLORS["muted"], align="center")

    def save(self):
        self.c.showPage()
        self.c.save()


def build_pdf():
    pdf = PDFDeck(PDF_OUT)

    pdf.rect(0, 0, 30, 21, COLORS["cream"])
    pdf.rect(0, 0, 30, 2.42, COLORS["sage"])
    pdf.rect(0, 2.42, 30, 0.08, COLORS["taupe_dark"])
    pdf.image_fit(UA_MATH_LOGO, 0.55, 0.48, 4.85, 1.18, border=False, fill=COLORS["sage"])
    pdf.image_fit(BIOSPHERE2_LOGO, 25.05, 0.68, 4.35, 0.85, border=False, fill=COLORS["sage"])
    pdf.text(5.25, 0.28, 19.5, 1.15,
             "Spontaneous Moss Colonization on Crushed Basalt:\n"
             "Imaging Pipeline, Color Reliability, and Topographic Controls on Persistence",
             size=23.5, bold=True, align="center", valign="middle", leading=25)
    pdf.text(7.4, 1.45, 15.2, 0.35, "Freddy Aguilar, Sophia Bui, Ava Mohr",
             size=17, align="center")
    pdf.text(7.4, 1.84, 15.2, 0.32, "The University of Arizona, Department of Mathematics",
             size=11.8, italic=True, align="center")

    left_x, left_w = 0.5, 6.35
    center_x, center_w = 7.15, 12.05
    right_x, right_w = 19.55, 9.95
    top_y = 2.8

    x, y, w, h = pdf.card(left_x, top_y, left_w, 3.05, "Why This System?")
    pdf.bullets(x, y, w, h, [
        "Biosphere 2's LEO Observatory is a controlled 30 m x 30 m hillslope built from crushed basalt.",
        "The substrate is a Mars-analog, soil-poor surface, yet moss and biological soil crusts colonized it without planting.",
        "The scientific question is whether this colonization is random or structured by hillslope geometry."
    ], size=9.3)
    x, y, w, h = pdf.card(left_x, 6.1, left_w, 2.15, "Research Questions", COLORS["sage_dark"])
    pdf.text(x + 0.05, y, w - 0.1, h,
             "Q1. Can overhead images reliably map moss area?\n\nQ2. Do topographic features predict where moss persists?",
             size=11.2, bold=True, align="center", valign="middle")
    x, y, w, h = pdf.card(left_x, 8.5, left_w, 2.75, "Experimental Setup")
    pdf.bullets(x, y, w, h, [
        "Sensor-rich laboratory system: 1,800+ devices, 16 sensor columns, and 5 cameras per slope.",
        "LiDAR DEM at 0.1 m resolution supports geometry-based analysis.",
        "Image archive includes RGB and RGB plus near-infrared imagery, filtered by quality."
    ], size=9.2)
    x, y, w, h = pdf.card(left_x, 11.5, left_w, 2.65, "Previous Work", COLORS["sage_dark"])
    pdf.bullets(x, y, w, h, [
        "Klausmeier (1999): slope water redistribution can generate vegetation stripes and spots.",
        "Rietkerk et al. (2002): plant-water-infiltration feedbacks produce non-uniform biological structure.",
        "This study tests that patterning logic on a controlled basalt hillslope."
    ], size=8.9)
    x, y, w, h = pdf.card(left_x, 14.4, left_w, 3.15, "Data and Time Window")
    pdf.bullets(x, y, w, h, [
        "Analysis dates: Nov 3, Nov 22, and Dec 7, 2022.",
        "Primary response: persistent moss mask from repeated positive detections.",
        "Topographic covariates: lateral distance, downstream position, trough depth, flow proxy, wall proximity, and soil thickness.",
        "Previous poster time frame: 2022-10-08 to 2024-01-01."
    ], size=8.85)
    x, y, w, h = pdf.card(left_x, 17.8, left_w, 2.25, "Poster Story", COLORS["taupe_dark"])
    pdf.text(x, y, w, h,
             "Act 1: a spontaneous Mars-analog colonization problem.\n"
             "Act 2: reliable imaging requires multi-date persistence, not single-date color.\n"
             "Act 3: topography predicts persistence through a bounded moisture optimum.",
             size=8.9)

    x, y, w, h = pdf.card(center_x, top_y, center_w, 2.8, "Methodology: Image-to-Moss Pipeline")
    stage_w = (w - 0.42) / 3
    for i, (num, head, desc) in enumerate([
        ("01", "Calibration", "Radial lens correction and per-camera rectification."),
        ("02", "Stitching", "Five corrected frames merged into one bath-coordinate mosaic."),
        ("03", "Persistence Proxy", "Broad RGB/HSV moss proxy retained only when positive on >=2 of 3 dates."),
    ]):
        sx = x + i * (stage_w + 0.21)
        pdf.rect(sx, y + 0.02, stage_w, h - 0.06, COLORS["cream"], COLORS["line"])
        pdf.text(sx + 0.08, y + 0.12, 0.48, 0.34, num, size=13, color=COLORS["red"], bold=True)
        pdf.text(sx + 0.62, y + 0.12, stage_w - 0.72, 0.34, head, size=10.8, color=COLORS["blue"], bold=True)
        pdf.text(sx + 0.1, y + 0.56, stage_w - 0.2, h - 0.66, desc, size=8.5)

    x, y, w, h = pdf.card(center_x, 5.9, center_w, 9.55, "Central Evidence Panels", COLORS["sage_dark"])
    img_y = y + 0.02
    center_slots = [
        (CENTER_IMAGES[0], x, img_y, 3.35, 7.65, img_y + 7.82),
        (CENTER_IMAGES[1], x + 3.58, img_y, 3.35, 5.15, img_y + 5.32),
        (CENTER_IMAGES[2], x + 7.16, img_y, 4.53, 5.55, img_y + 5.72),
    ]
    for (path, head, desc), sx, sy, sw, sh, cy in center_slots:
        pdf.image_fit(path, sx, sy, sw, sh)
        pdf.text(sx, cy, sw, 0.28, head, size=9.1, color=COLORS["blue"], bold=True, align="center")
        pdf.text(sx + 0.05, cy + 0.31, sw - 0.1, 0.52, desc, size=7.0, color=COLORS["muted"], align="center")
    pdf.rect(x, y + 8.72, w, 0.62, COLORS["taupe_dark"])
    pdf.text(x + 0.16, y + 8.8, w - 0.32, 0.45,
             "Stable moss is concentrated on off-center shoulder zones, not in the deepest drainage corridor.",
             size=10.4, color=RGBColor(255, 255, 255), bold=True, align="center", valign="middle")

    x, y, w, h = pdf.card(center_x, 15.75, center_w, 4.3, "Findings 01-02: Persistence Beats Color")
    pdf.text(x, y, w * 0.58, h - 0.1,
             "Finding 01 - Moss is spatially persistent.\n"
             "The same regions remain moss-positive in November and December. Later greening is mostly color intensification in already occupied regions.\n\n"
             "Finding 02 - Color thresholding is unreliable.\n"
             "Hydration and irrigation can shift moss from yellow-green to vivid green within hours, so single-date green area can mimic growth.",
             size=8.7)
    stat_x = x + w * 0.61
    pdf.stat(stat_x, y + 0.04, 1.55, 1.05, "26.6%", "stable on >=2 dates")
    pdf.stat(stat_x + 1.75, y + 0.04, 1.55, 1.05, "7.1%", "strong on all 3 dates")
    pdf.stat(stat_x, y + 1.33, 1.55, 1.05, "3.7 m", "median lateral offset")
    pdf.stat(stat_x + 1.75, y + 1.33, 1.55, 1.05, "0.088 m", "median trough depth")
    pdf.text(stat_x, y + 2.66, 3.28, 0.82,
             "Interpretation: the response variable is persistence, not raw greenness.",
             size=8.2, color=COLORS["taupe_dark"], bold=True, align="center", valign="middle")

    x, y, w, h = pdf.card(right_x, top_y, right_w, 2.9, "Finding 03: Topographic Control", COLORS["taupe_dark"])
    pdf.rect(x, y + 0.02, w, 0.72, COLORS["blue"])
    pdf.text(x + 0.12, y + 0.12, w - 0.24, 0.5,
             "The wettest corridor is not the most favorable habitat.",
             size=12.1, color=RGBColor(255, 255, 255), bold=True, align="center", valign="middle")
    pdf.bullets(x, y + 0.92, w, h - 1.0, [
        "Persistent moss increases away from the centerline.",
        "Trough-depth response is non-monotone, peaking near an intermediate depth.",
        "Deep, high-flow drainage corridors behave as exclusion zones."
    ], size=8.9)

    x, y, w, h = pdf.card(right_x, 6.0, right_w, 5.05, "Quantitative Evidence")
    pdf.image_fit(ARTICLE_RESPONSE_CURVES, x, y + 0.02, w * 0.54, h - 0.34)
    pdf.image_fit(ARTICLE_CALIBRATION, x + w * 0.57, y + 0.02, w * 0.43, h - 0.34)
    pdf.text(x, y + h - 0.27, w, 0.22,
             "Response curves and suitability-decile calibration from the technical note.",
             size=7.0, color=COLORS["muted"], align="center")

    x, y, w, h = pdf.card(right_x, 11.35, right_w, 3.55, "Logistic Suitability Ranking Model", COLORS["sage_dark"])
    pdf.rect(x, y + 0.04, w, 0.57, COLORS["cream"], COLORS["line"])
    pdf.text(x + 0.08, y + 0.16, w - 0.16, 0.32,
             "logit P(M=1) = beta0 + beta_w W - beta_f F - gamma(d-d0)^2 + beta_t T",
             size=8.2, bold=True, align="center")
    pdf.text(x, y + 0.82, w * 0.51, 1.74,
             "Predictors\n"
             "W: wall proximity (+)\n"
             "F: log(1+flow) (-)\n"
             "d: trough depth, optimum d0=0.407 m\n"
             "T: soil thickness (+)",
             size=8.3)
    pdf.text(x + w * 0.55, y + 0.82, w * 0.42, 1.74,
             "Performance\n"
             "Training AUC: 0.724\n"
             "Spatial CV AUC: 0.657\n\n"
             "Suitability deciles rank stable-moss fraction from about 7% to 48%.",
             size=8.3)
    pdf.text(x, y + 2.72, w, 0.34,
             "Use as first-order ranking, not a complete mechanistic law.",
             size=8.2, color=COLORS["taupe_dark"], bold=True, align="center")

    x, y, w, h = pdf.card(right_x, 15.2, right_w, 2.35, "Why the Trough May Suppress Moss")
    pdf.bullets(x, y, w, h, [
        "excess saturation and low oxygen",
        "surface flushing of spores and fines",
        "sediment redistribution",
        "mechanical surface instability",
        "moisture-window limitation"
    ], size=8.3)
    x, y, w, h = pdf.card(right_x, 17.85, right_w, 2.2, "Conclusions and Next Steps", COLORS["taupe_dark"])
    pdf.text(x, y, w * 0.54, h - 0.05,
             "Q1: Yes - use multi-date persistence to suppress hydration artifacts.\n\n"
             "Q2: Yes - topography ranks persistence, with shoulder zones favored over active drainage corridors.\n\n"
             "Implication: pioneer colonizers occupy off-channel shoulders on Mars-analog basalt.",
             size=7.8)
    pdf.text(x + w * 0.58, y, w * 0.4, h - 0.05,
             "Next measurements\n"
             "- soil moisture time series\n"
             "- saturation duration\n"
             "- tracer residence time\n"
             "- CO2 flux vs moss density",
             size=8.1)

    pdf.text(0.55, 20.35, 13.4, 0.28,
             "Scientific source: moss_technical_note_scientific_revision.pdf; poster narrative: Findings-content.md.md",
             size=6.5, color=COLORS["muted"])
    pdf.text(20.0, 20.35, 9.45, 0.28,
             "References: Klausmeier 1999; Rietkerk et al. 2002; LEO Observatory / Biosphere 2 context.",
             size=6.5, color=COLORS["muted"], align="right")

    pdf.save()
    print(PDF_OUT)


if __name__ == "__main__":
    build()
